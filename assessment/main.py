import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import openai


class PresentationGenerator:
    def __init__(self, pixabay_api_key, gpt_api_key):
        self.gpt_api_key = gpt_api_key

    def generate_content(self, topic):
        openai.api_key = self.gpt_api_key
        prompt = f"Write about {topic}"
        response = openai.Completion.create(
            engine="davinci", prompt=prompt, max_tokens=30
        )
        content = response.choices[0].text.strip()
        return content

    def add_line_break(self, text):
        new_input = ""
        for i, letter in enumerate(text):
            if i % 30== 0:
                new_input += '\n'
            new_input += letter
        new_input = new_input[1:]
        return new_input

    def add_text(self, slide, text, left, top, width, height, font_size, is_font_bold, font_color):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = text
        p.font.bold = is_font_bold
        p.font.size = font_size
        p.font.color.rgb = font_color
        return

    def generate_slide(self, slide_data, prs):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = slide_data["title"]["props"]["left"]
        top = slide_data["title"]["props"]["top"]
        height = slide_data["title"]["props"]["height"]
        width = slide_data["title"]["props"]["width"]
        font_size = slide_data["title"]["props"]["font_size"]
        font_color = RGBColor(0, 0, 139)
        self.add_text(slide, slide_data["title"]["value"], left, top, width, height, font_size, True, font_color)

        if slide_data.get("paragraphs"):
            stack_type = slide_data["paragraphs"]["stack_type"]
            if not stack_type:
                values = slide_data["paragraphs"]["values"]
                for value in values:
                    text = value["text"]
                    if not text:
                        text = self.generate_content(value["ai_input"])
                        text = self.add_line_break(text)
                    p_left = value["props"]["left"]
                    p_top = value["props"]["top"]
                    p_height = value["props"]["height"]
                    p_width = value["props"]["width"]
                    p_font_size = value["props"]["font_size"]
                    p_font_color = RGBColor(0, 0, 0)
                    self.add_text(slide, text, p_left, p_top, p_width, p_height, p_font_size, False,
                                  p_font_color)

        path = slide_data["image"]["path"]
        image_left = slide_data["image"]["props"]["left"]
        image_top = slide_data["image"]["props"]["top"]
        image_height = slide_data["image"]["props"]["height"]
        slide.shapes.add_picture(path, image_left, image_top, height=image_height)

        return prs

    def generate_ppt(self, slides):
        prs = Presentation()
        for slide_data in slides:
            prs = self.generate_slide(slide_data, prs)
        prs.save("presentation_final.pptx")


# Usage
pixabay_api_key = "39734621-3c7748887bfc057290f4ce087"
gpt_api_key = "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"

slides = [{"title": {"value": "Presentation Skills",
                     "props": {"left": Inches(0.8), "top": Inches(3), "height": Inches(1), "width": Inches(1),
                               "font_size": Pt(35)}}, "is_content_required": False, "body": "",
           "paragraphs": {},
           "image": {"path": "media/picture1.jpg",
                     "props": {"left": Inches(5), "top": Inches(1), "height": Inches(5)}}},
          {"title": {"value": "Introduction",
                     "props": {"left": Inches(3.8), "top": Inches(0.5), "height": Inches(1), "width": Inches(1),
                               "font_size": Pt(20)}}, "is_content_required": False, "body": "",
           "paragraphs": {"stack_type": "", "values": [{
               "text": "",
               "ai_input": "introduction to powerpoint presentation skills in 10 words",
               "props": {"left": Inches(1.5), "top": Inches(5.0),
                         "height": Inches(1), "width": Inches(1),
                         "font_size": Pt(10)}},
               {
                   "text": "",
                   "ai_input": "introduction to technical skills required in powerpoint presentation in 10 words",
                   "props": {"left": Inches(4.8), "top": Inches(5.0),
                             "height": Inches(1), "width": Inches(1),
                             "font_size": Pt(10)}},
               {
                   "text": "",
                   "ai_input": "introduction to  analytical skills required in powerpoint presentation skills in 10 words",
                   "props": {"left": Inches(7.0), "top": Inches(5.0),
                             "height": Inches(1), "width": Inches(1),
                             "font_size": Pt(10)}}
           ]},
           "image": {"path": "media/picture6.png",
                     "props": {"left": Inches(1.2), "top": Inches(1.5), "height": Inches(3.5)}}},

          {"title": {"value": "The Four Types of Presentation",
                     "props": {"left": Inches(0.8), "top": Inches(2), "height": Inches(1), "width": Inches(1),
                               "font_size": Pt(20)}}, "is_content_required": False, "body": "",
           "paragraphs": {"stack_type": "", "values": [{
               "text": "There are four kinds of presentations:\n informational, instructional ,\nstimulating, and convincing.",
               "props": {"left": Inches(0.8), "top": Inches(3.5),
                         "height": Inches(1), "width": Inches(1),
                         "font_size": Pt(17)}}]},
           "image": {"path": "media/picture2.jpg",
                     "props": {"left": Inches(5), "top": Inches(1), "height": Inches(5)}}},
          {"title": {"value": "Course Evaluation",
                     "props": {"left": Inches(0.8), "top": Inches(2), "height": Inches(1), "width": Inches(1),
                               "font_size": Pt(20)}}, "is_content_required": False, "body": "",
           "paragraphs": {"stack_type": "", "values": [{
               "text": "Kindly spare a few moments to fill out the \n course feedback form.",
               "props": {"left": Inches(0.8), "top": Inches(3.5),
                         "height": Inches(1), "width": Inches(1),
                         "font_size": Pt(17)}}]},
           "image": {"path": "media/picture5.png",
                     "props": {"left": Inches(5.4), "top": Inches(1), "height": Inches(5)}}},
          ]

presentation_generator = PresentationGenerator(pixabay_api_key, gpt_api_key)
presentation_generator.generate_ppt(slides)
